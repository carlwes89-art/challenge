from pathlib import Path

from pypdf import PdfReader
from docx import Document
from pptx import Presentation
from bs4 import BeautifulSoup
import pandas as pd


# ============================================================
# PDF
# ============================================================

def load_pdf(path):

    reader = PdfReader(path)

    documents = []

    for page_number, page in enumerate(reader.pages):

        text = page.extract_text()

        if text and text.strip():

            documents.append({
                "text": text.strip(),

                "metadata": {
                    "source": str(path),
                    "type": "pdf",
                    "page": page_number + 1
                }
            })

    return documents


# ============================================================
# DOCX
# ============================================================

def load_docx(path):

    document = Document(path)

    paragraphs = []

    for paragraph in document.paragraphs:

        text = paragraph.text.strip()

        if text:
            paragraphs.append(text)

    text = "\n".join(paragraphs)

    if not text:
        return []

    return [{
        "text": text,

        "metadata": {
            "source": str(path),
            "type": "docx"
        }
    }]


# ============================================================
# TXT / MD
# ============================================================

def load_text(path):

    with open(
        path,
        "r",
        encoding="utf-8"
    ) as file:

        text = file.read()

    if not text.strip():
        return []

    return [{
        "text": text.strip(),

        "metadata": {
            "source": str(path),
            "type": path.suffix.lower().replace(".", "")
        }
    }]


# ============================================================
# CSV
# ============================================================

def load_csv(path):

    df = pd.read_csv(path)

    documents = []

    for index, row in df.iterrows():

        values = []

        for column in df.columns:

            value = row[column]

            if pd.notna(value):

                values.append(
                    f"{column}: {value}"
                )

        text = "\n".join(values)

        if text.strip():

            documents.append({
                "text": text,

                "metadata": {
                    "source": str(path),
                    "type": "csv",
                    "row": index + 1
                }
            })

    return documents


# ============================================================
# XLSX
# ============================================================

def load_xlsx(path):

    sheets = pd.read_excel(
        path,
        sheet_name=None
    )

    documents = []

    for sheet_name, df in sheets.items():

        for index, row in df.iterrows():

            values = []

            for column in df.columns:

                value = row[column]

                if pd.notna(value):

                    values.append(
                        f"{column}: {value}"
                    )

            text = "\n".join(values)

            if text.strip():

                documents.append({
                    "text": text,

                    "metadata": {
                        "source": str(path),
                        "type": "xlsx",
                        "sheet": sheet_name,
                        "row": index + 1
                    }
                })

    return documents


# ============================================================
# PPTX
# ============================================================

def load_pptx(path):

    presentation = Presentation(path)

    documents = []

    for slide_number, slide in enumerate(
        presentation.slides
    ):

        texts = []

        for shape in slide.shapes:

            if hasattr(shape, "text"):

                text = shape.text.strip()

                if text:
                    texts.append(text)

        if texts:

            documents.append({
                "text": "\n".join(texts),

                "metadata": {
                    "source": str(path),
                    "type": "pptx",
                    "slide": slide_number + 1
                }
            })

    return documents


# ============================================================
# HTML
# ============================================================

def load_html(path):

    with open(
        path,
        "r",
        encoding="utf-8"
    ) as file:

        html = file.read()

    soup = BeautifulSoup(
        html,
        "html.parser"
    )

    # Supprime les éléments inutiles
    for element in soup(
        ["script", "style", "noscript"]
    ):
        element.decompose()

    text = soup.get_text(
        separator="\n"
    )

    if not text.strip():
        return []

    return [{
        "text": text.strip(),

        "metadata": {
            "source": str(path),
            "type": "html"
        }
    }]


# ============================================================
# LOADER PRINCIPAL
# ============================================================

def load_document(path):

    path = Path(path)

    extension = path.suffix.lower()

    if extension == ".pdf":
        return load_pdf(path)

    elif extension == ".docx":
        return load_docx(path)

    elif extension in [".txt", ".md"]:
        return load_text(path)

    elif extension == ".csv":
        return load_csv(path)

    elif extension == ".xlsx":
        return load_xlsx(path)

    elif extension == ".pptx":
        return load_pptx(path)

    elif extension in [".html", ".htm"]:
        return load_html(path)

    else:

        raise ValueError(
            f"Format non supporté : {extension}"
        )


# ============================================================
# CHARGER TOUS LES DOCUMENTS D'UN DOSSIER
# ============================================================

SUPPORTED_EXTENSIONS = {
    ".pdf",
    ".docx",
    ".txt",
    ".md",
    ".csv",
    ".xlsx",
    ".pptx",
    ".html",
    ".htm"
}


def load_directory(directory):

    directory = Path(directory)

    all_documents = []

    for path in directory.rglob("*"):

        if not path.is_file():
            continue

        if path.suffix.lower() not in SUPPORTED_EXTENSIONS:
            continue

        try:

            documents = load_document(path)

            all_documents.extend(documents)

            print(
                f"[OK] {path} "
                f"→ {len(documents)} éléments"
            )

        except Exception as error:

            print(
                f"[ERREUR] {path} : {error}"
            )

    return all_documents